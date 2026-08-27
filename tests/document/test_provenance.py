"""Provenance, verified by OPENING the files. Spec 082, FR-006..FR-010, FR-034.

SC-009 says every generated document states its generation time and that NetClaw
produced it, "verified by opening the file, not by inspecting code". That wording exists
because spec 080 shipped `fgt_system_status` returning three nulls past 24 passing
tests: the tests asserted the envelope was well-formed, and it was — the payload was
half empty.

So every check below reparses the artefact with the same library a reader's application
would use, and asserts on what is inside it.
"""

from __future__ import annotations

import json
import os
import zipfile

from _harness import FAILURES, abspath, check, cleanup, run, sandbox  # noqa: F401

import envelope  # noqa: E402
import output  # noqa: E402
from outcomes import Outcome  # noqa: E402
from provenance import DocumentStamp, SourceLedger  # noqa: E402
from writers import docx_writer, pptx_writer, xlsx_writer  # noqa: E402

DOCX_PAYLOAD = {
    "title": "Change Record CHG0012345",
    "blocks": [
        {"type": "heading", "text": "Pre-change state", "level": 1},
        {"type": "keyvalue", "pairs": [
            {"label": "Hostname", "value": {"v": "fgt-01", "src": "fgt_system_status",
                                            "as_of": "2026-08-01T09:00:00Z"}},
            {"label": "Serial", "value": {"unavailable": "device did not answer"}},
        ]},
        {"type": "table", "caption": "Interfaces",
         "columns": ["Interface", "Admin state", "Oper state"],
         "rows": [[{"v": "port1", "src": "fgt_list_interfaces", "device": "fgt-01"},
                   {"v": "up", "src": "fgt_list_interfaces", "device": "fgt-01"},
                   {"v": "down", "src": "fgt_list_interfaces", "device": "fgt-01"}]]},
    ],
}

XLSX_PAYLOAD = {"sheets": [{
    "name": "Interfaces",
    "columns": ["Device", "Interface", "Admin state", "Oper state"],
    "rows": [[{"v": "fgt-01", "src": "fgt_list_interfaces", "device": "fgt-01",
               "as_of": "2026-08-01T09:00:00Z"},
              {"v": "port1", "src": "fgt_list_interfaces", "device": "fgt-01"},
              {"v": "up", "src": "fgt_list_interfaces", "device": "fgt-01"},
              {"v": "down", "src": "fgt_list_interfaces", "device": "fgt-01"}]],
}]}

PPTX_PAYLOAD = {"title": "Posture review", "slides": [
    {"layout": "figure", "title": "Scope", "figures": [
        {"label": "Interfaces", "value": {"v": 2, "src": "fgt_list_interfaces",
                                          "device": "fgt-01", "as_of": "2026-08-01T09:00:00Z"}}]},
]}


def _emit(kind, tool, payload, builder):
    stamp = DocumentStamp(tool=tool)
    ledger = SourceLedger()
    data, caveats = builder(payload, stamp, ledger)
    path, suffix = output.reserve(kind, "provenance")
    path.write_bytes(data)
    art = output.finalize(path, suffix)
    resp = envelope.emit(tool=tool, outcome=Outcome.OK, artifact=art.as_dict(),
                         ledger=ledger, stamp=stamp, caveats=caveats)
    return path, resp, stamp


def _docx_text(path):
    from docx import Document

    d = Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for r in t.rows:
            parts.append(" | ".join(c.text for c in r.cells))
    parts.append(d.sections[0].footer.paragraphs[0].text)
    return "\n".join(parts), d


def test_docx_carries_the_stamp_and_sources():
    d = sandbox()
    try:
        path, resp, stamp = _emit("docx", "docx_write", DOCX_PAYLOAD, docx_writer.build)
        text, doc = _docx_text(path)

        footer = doc.sections[0].footer.paragraphs[0].text
        check("the footer names NetClaw", "NetClaw" in footer, footer)
        check("the footer carries a generation time", stamp.generated_at in footer, footer)
        check("the footer appears on every page (section footer)", bool(footer.strip()))

        check("a Sources section exists", "Sources" in text, "no Sources heading")
        for src in ("fgt_system_status", "fgt_list_interfaces"):
            check(f"the Sources section lists {src}", src in text, "missing from the document")

        headers = [c.text for t in doc.tables for c in t.rows[0].cells]
        check("tables carry a Source column", "Source" in headers, str(headers))
        check("tables carry an As of column", "As of" in headers, str(headers))

        check(
            "the source's own as-of is in the document",
            "2026-08-01T09:00:00Z" in text,
            "the source's collection time was lost",
        )
        check(
            "the source's as-of is distinguishable from the generation time",
            "2026-08-01T09:00:00Z" != stamp.generated_at,
            "they are the same value, so FR-010 is untestable here",
        )

        for t in doc.tables:
            headers_t = [c.text for c in t.rows[0].cells]
            if "Source" not in headers_t:
                continue
            col = headers_t.index("Source")
            for row in t.rows[1:]:
                cell = row.cells[col].text.strip()
                check(
                    f"no data row has an empty Source cell (row starts {row.cells[0].text[:18]!r})",
                    bool(cell),
                    "an unattributed row shipped",
                )
    finally:
        cleanup(d)


def test_xlsx_carries_the_stamp_and_sources():
    d = sandbox()
    try:
        import openpyxl

        path, resp, stamp = _emit("xlsx", "xlsx_write", XLSX_PAYLOAD, xlsx_writer.build)
        wb = openpyxl.load_workbook(str(path))

        check("a Sources sheet exists", "Sources" in wb.sheetnames, str(wb.sheetnames))
        ws = wb["Interfaces"]
        banner = str(ws.cell(row=1, column=1).value or "")
        check("the banner names NetClaw", "NetClaw" in banner, banner[:80])
        check("the banner carries a generation time", stamp.generated_at in banner, banner[:80])

        header_row = next(r for r in range(1, 6)
                          if ws.cell(row=r, column=1).value == "Device")
        headers = [ws.cell(row=header_row, column=c).value for c in range(1, 8)]
        check("a Source column was appended", "Source" in headers, str(headers))
        check("an As of column was appended", "As of" in headers, str(headers))

        src_col = headers.index("Source") + 1
        val = ws.cell(row=header_row + 1, column=src_col).value
        check("the data row carries its source", val and "fgt_list_interfaces" in str(val), str(val))
        age_col = headers.index("As of") + 1
        check(
            "the source's own as-of is in the sheet",
            str(ws.cell(row=header_row + 1, column=age_col).value) == "2026-08-01T09:00:00Z",
            str(ws.cell(row=header_row + 1, column=age_col).value),
        )

        srcsheet = wb["Sources"]
        listed = [srcsheet.cell(row=r, column=1).value for r in range(1, 10)]
        check("the Sources sheet lists the tool", "fgt_list_interfaces" in listed, str(listed))
    finally:
        cleanup(d)


def test_pptx_source_is_in_a_shape_not_only_notes():
    """SC-010b. Speaker notes are invisible in presentation and print — if attribution
    lives only there, it is not attribution."""
    d = sandbox()
    try:
        from pptx import Presentation

        path, resp, stamp = _emit("pptx", "pptx_write", PPTX_PAYLOAD, pptx_writer.build)
        prs = Presentation(str(path))

        check("a Sources slide exists",
              any(s.shapes.title and s.shapes.title.text == "Sources" for s in prs.slides),
              "no Sources slide")

        for i, slide in enumerate(prs.slides):
            shape_text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
            check(
                f"slide {i} carries the stamp in a SHAPE",
                "NetClaw" in shape_text,
                "attribution is not present in any visible shape on this slide",
            )

        figure_slide = next(s for s in prs.slides
                            if s.shapes.title and s.shapes.title.text == "Scope")
        shape_text = "\n".join(sh.text_frame.text for sh in figure_slide.shapes if sh.has_text_frame)
        notes_text = figure_slide.notes_slide.notes_text_frame.text
        check("the figure's source is in a visible shape", "fgt_list_interfaces" in shape_text,
              "the source is not on the slide")
        check("notes are additive, not the mechanism", "fgt_list_interfaces" in notes_text,
              "notes were expected to repeat the detail")
    finally:
        cleanup(d)


def test_a_caller_cannot_omit_provenance():
    """SC-023. The guarantee is structural, not a convention."""
    d = sandbox()
    try:
        raised = None
        try:
            envelope.emit(tool="docx_write", outcome=Outcome.OK,
                          artifact={"path": "x", "bytes": 1, "created_at": "now"},
                          ledger=SourceLedger(), stamp=DocumentStamp(tool="docx_write"))
        except envelope.ProvenanceError as exc:
            raised = exc
        check("emitting an artefact with an empty ledger raises", raised is not None,
              "an unattributed document was emittable")

        resp = envelope.emit(tool="docx_write", outcome=Outcome.OK, data={"x": 1},
                             ledger=SourceLedger(), stamp=DocumentStamp(tool="docx_write"))
        check("a non-artefact response still carries generated_at", bool(resp["generated_at"]))
        check("a non-artefact response still carries generated_by", bool(resp["generated_by"]))

        stamp = DocumentStamp(tool="docx_write")
        check("a caller cannot choose generated_by",
              stamp.generated_by.startswith("NetClaw document-mcp"), stamp.generated_by)
    finally:
        cleanup(d)


def test_gait_records_successes_and_refusals():
    """FR-034 / SC-019. The defect /speckit.analyze caught in specs 076 and 080 was GAIT
    having verification but no implementation. A refusal must be audited too."""
    d = sandbox()
    try:
        log = os.environ["DOCUMENT_AUDIT_LOG"]
        _emit("docx", "docx_write", DOCX_PAYLOAD, docx_writer.build)
        envelope.refused(tool="docx_write", reason="template supplied",
                         outcome=Outcome.REFUSED_TEMPLATE)

        check("the audit log was created", os.path.exists(log), log)
        records = [json.loads(line) for line in open(log, encoding="utf-8") if line.strip()]
        check("at least two records were written", len(records) >= 2, str(len(records)))
        outcomes = {r["outcome"] for r in records}
        check("the success is audited", "ok" in outcomes or "written_with_gaps" in outcomes, str(outcomes))
        check("the REFUSAL is audited", "refused_template" in outcomes, str(outcomes))
        for r in records:
            check(f"record for {r['tool']} names the component",
                  r["component"] == "document-mcp", str(r))
            check(f"record for {r['tool']}/{r['outcome']} carries a timestamp",
                  bool(r.get("ts")), str(r))
    finally:
        cleanup(d)


def test_provenance_survives_the_forwarding_path():
    """SC-010b. Copy the table out and print the file — does attribution survive?

    Extracting the raw XML text is the copy-paste path; asserting the attribution is
    absent from docProps is the "not hidden in metadata" half.
    """
    d = sandbox()
    try:
        path, _resp, _stamp = _emit("docx", "docx_write", DOCX_PAYLOAD, docx_writer.build)
        with zipfile.ZipFile(str(path)) as z:
            body = z.read("word/document.xml").decode()
        check(
            "attribution is in the document body, not only in metadata",
            "fgt_list_interfaces" in body,
            "the source name is absent from word/document.xml",
        )

        from docx import Document

        doc = Document(str(path))
        table = next(t for t in doc.tables if "Interface" in [c.text for c in t.rows[0].cells])
        copied = "\n".join(" ".join(c.text for c in r.cells) for r in table.rows)
        check(
            "a copied table still carries its source",
            "fgt_list_interfaces" in copied,
            "attribution was lost on copy — it must not live in a comment or tooltip",
        )

        with zipfile.ZipFile(str(path)) as z:
            names = z.namelist()
        check(
            "no Word comment part carries the provenance",
            "word/comments.xml" not in names,
            "provenance may be additive in comments, but must not depend on them",
        )
    finally:
        cleanup(d)


TESTS = [
    test_docx_carries_the_stamp_and_sources,
    test_xlsx_carries_the_stamp_and_sources,
    test_pptx_source_is_in_a_shape_not_only_notes,
    test_a_caller_cannot_omit_provenance,
    test_gait_records_successes_and_refusals,
    test_provenance_survives_the_forwarding_path,
]

if __name__ == "__main__":
    raise SystemExit(run(TESTS, "provenance"))
