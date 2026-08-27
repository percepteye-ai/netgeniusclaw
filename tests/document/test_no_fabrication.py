"""Never fabricate to fill a blank. Spec 082, FR-001..FR-005, FR-015, FR-023a, FR-027a/b.

This is the suite that matters most, and every assertion in it reopens the produced
file. The failure it guards against is not a crash — it is a professional-looking
document that reads as complete and is not.
"""

from __future__ import annotations

from _harness import FAILURES, check, cleanup, run, sandbox  # noqa: F401

import output  # noqa: E402
from outcomes import Outcome, ValueRefused  # noqa: E402
from provenance import DocumentStamp, SourceLedger  # noqa: E402
from writers import docx_writer, pptx_writer, xlsx_writer  # noqa: E402

# Every string a gap must never be mistaken for.
LAUNDERED = ["N/A", "n/a", "TBD", "Unknown", "None", "--"]


def _build(builder, payload, tool):
    stamp = DocumentStamp(tool=tool)
    ledger = SourceLedger()
    data, caveats = builder(payload, stamp, ledger)
    path, suffix = output.reserve({"docx_write": "docx", "xlsx_write": "xlsx",
                                   "pptx_write": "pptx"}[tool], "nofab")
    path.write_bytes(data)
    return path, caveats, stamp, ledger


def _docx_text(path):
    from docx import Document

    d = Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for t in d.tables:
        for r in t.rows:
            parts.append(" | ".join(c.text for c in r.cells))
    return "\n".join(parts)


def test_gaps_render_explicitly_in_the_document():
    d = sandbox()
    try:
        payload = {"title": "T", "blocks": [{"type": "keyvalue", "pairs": [
            {"label": "Hostname", "value": {"v": "fgt-01", "src": "fgt_system_status"}},
            {"label": "Serial", "value": {"unavailable": "device did not answer"}},
            {"label": "Uptime", "value": {"failed": "connection refused"}},
            {"label": "Banner", "value": {"v": "", "src": "fgt_system_status"}},
        ]}]}
        path, _c, _s, _l = _build(docx_writer.build, payload, "docx_write")
        text = _docx_text(path)

        check("an unavailable field says NOT AVAILABLE", "NOT AVAILABLE" in text)
        check("its reason is in the document", "device did not answer" in text)
        check("a failed field says RETRIEVAL FAILED", "RETRIEVAL FAILED" in text)
        check("its reason is in the document", "connection refused" in text)
        check("a genuinely-empty source value says (empty)", "(empty)" in text)
        check(
            "unavailable and failed are textually distinct",
            "NOT AVAILABLE" in text and "RETRIEVAL FAILED" in text,
            "a dead query and an empty result must not read the same",
        )
        for token in LAUNDERED:
            check(f"the document contains no {token!r} placeholder", token not in text,
                  "a plausible placeholder is how a guess gets laundered into a record")
    finally:
        cleanup(d)


def test_a_gapped_document_cannot_report_as_clean():
    d = sandbox()
    try:
        import envelope

        payload = {"title": "T", "blocks": [{"type": "keyvalue", "pairs": [
            {"label": "A", "value": {"v": 1, "src": "x"}},
            {"label": "B", "value": {"unavailable": "no data"}},
        ]}]}
        path, caveats, stamp, ledger = _build(docx_writer.build, payload, "docx_write")
        art = output.finalize(path, None)
        resp = envelope.emit(tool="docx_write", outcome=Outcome.OK, artifact=art.as_dict(),
                             ledger=ledger, stamp=stamp, caveats=caveats)
        check("the caller passed ok and got written_with_gaps",
              resp["outcome"] == Outcome.WRITTEN_WITH_GAPS.value, resp["outcome"])
        check("the gap count is reported", resp["gaps"]["unavailable"] == 1, str(resp["gaps"]))
        check("a caveat states the document is incomplete",
              any("incomplete" in c for c in resp["caveats"]), str(resp["caveats"]))

        clean = {"title": "T", "blocks": [{"type": "keyvalue", "pairs": [
            {"label": "A", "value": {"v": 1, "src": "x"}}]}]}
        p2, c2, s2, l2 = _build(docx_writer.build, clean, "docx_write")
        r2 = envelope.emit(tool="docx_write", outcome=Outcome.OK,
                           artifact=output.finalize(p2, None).as_dict(), ledger=l2, stamp=s2,
                           caveats=c2)
        check("a genuinely complete document reports ok", r2["outcome"] == "ok", r2["outcome"])
    finally:
        cleanup(d)


def test_failed_devices_are_rows_not_omissions():
    """SC-004. A shorter spreadsheet reads as a smaller estate."""
    d = sandbox()
    try:
        import openpyxl

        payload = {"sheets": [{
            "name": "Interfaces",
            "columns": ["Device", "Interface", "Admin state", "Oper state"],
            "rows": [[{"v": "fgt-01", "src": "t"}, {"v": "port1", "src": "t"},
                      {"v": "up", "src": "t"}, {"v": "down", "src": "t"}]],
            "failed_rows": [{"label": "fgt-02", "failed": "connection refused"},
                            {"label": "fgt-03", "failed": "auth expired"}],
        }]}
        path, caveats, _s, _l = _build(xlsx_writer.build, payload, "xlsx_write")
        ws = openpyxl.load_workbook(str(path))["Interfaces"]

        cells = [str(ws.cell(row=r, column=c).value or "")
                 for r in range(1, 12) for c in range(1, 7)]
        blob = "\n".join(cells)
        check("fgt-02 appears in the sheet", "fgt-02" in blob, "a failed device was omitted")
        check("fgt-03 appears in the sheet", "fgt-03" in blob, "a failed device was omitted")
        check("failed rows are marked as failed", blob.count("RETRIEVAL FAILED") >= 2, blob[:200])

        banner = str(ws.cell(row=1, column=1).value or "")
        check("the banner reports attempted", "3 attempted" in banner, banner[:90])
        check("the banner reports the failures", "2 failed" in banner, banner[:90])
        check("the banner distinguishes returned data", "1 returned data" in banner, banner[:90])
        check("a caveat names the unreached devices",
              any("could not be reached" in c for c in caveats), str(caveats))

        check("a failed row with no reason is refused",
              _refused(xlsx_writer.build, {"sheets": [{"name": "S", "columns": ["A"], "rows": [],
                                                       "failed_rows": [{"label": "x"}]}]}) is not None)
    finally:
        cleanup(d)


def _refused(builder, payload):
    try:
        builder(payload, DocumentStamp(tool="t"), SourceLedger())
    except ValueRefused as exc:
        return exc
    return None


def test_admin_and_oper_state_stay_separate():
    """FR-015 / SC-006 — the distinction spec 080's completion established."""
    d = sandbox()
    try:
        exc = _refused(xlsx_writer.build,
                       {"sheets": [{"name": "S", "columns": ["Device", "Status"], "rows": []}]})
        check("a merged 'Status' column is refused", exc is not None, "it was accepted")
        if exc:
            check("the refusal is typed", exc.outcome is Outcome.REFUSED_MERGED_STATUS, str(exc.outcome))
            check("the refusal explains the distinction",
                  "administratively up" in str(exc), str(exc)[:120])
        for header in ("state", "Up/Down", "link status"):
            check(f"a merged {header!r} column is refused",
                  _refused(xlsx_writer.build,
                           {"sheets": [{"name": "S", "columns": ["Device", header],
                                        "rows": []}]}) is not None)
        check(
            "a summary column alongside both real columns is allowed",
            _refused(xlsx_writer.build,
                     {"sheets": [{"name": "S",
                                  "columns": ["Admin state", "Oper state", "Status"],
                                  "rows": []}]}) is None,
            "a derived summary is legitimate when both facts are also present",
        )
    finally:
        cleanup(d)


def test_all_three_office_formats_refuse_a_template():
    """SC-008b. Parametrised so a fourth format cannot be added without it."""
    cases = [
        ("docx_write", docx_writer.build, {"title": "T", "blocks": [], "template": "corp.docx"}),
        ("xlsx_write", xlsx_writer.build,
         {"sheets": [{"name": "S", "columns": ["A"], "rows": []}], "template": "corp.xlsx"}),
        ("pptx_write", pptx_writer.build, {"title": "T", "slides": [], "template": "corp.pptx"}),
    ]
    for tool, builder, payload in cases:
        exc = _refused(builder, payload)
        check(f"{tool} refuses a template", exc is not None,
              "the template was silently ignored — the operator gets an unbranded "
              "document believing it is branded")
        if exc:
            check(f"{tool}'s refusal is typed", exc.outcome is Outcome.REFUSED_TEMPLATE, str(exc.outcome))
            check(f"{tool}'s refusal explains why", "fabrication pressure" in str(exc), str(exc)[:100])
    for alias in ("template_path", "base_document"):
        check(f"the alias {alias!r} is also refused",
              _refused(docx_writer.build, {"title": "T", "blocks": [], alias: "x"}) is not None)


def test_bounds_are_stated_inside_the_document():
    """FR-027a / SC-026 — for blocks and slides, not only worksheet rows."""
    import os

    d = sandbox()
    try:
        os.environ["DOCUMENT_MAX_BLOCKS"] = "3"
        os.environ["DOCUMENT_MAX_SLIDES"] = "2"
        os.environ["DOCUMENT_MAX_ROWS"] = "2"

        blocks = [{"type": "paragraph", "text": f"para {chr(97 + i)}"} for i in range(10)]
        path, _c, stamp, _l = _build(docx_writer.build, {"title": "T", "blocks": blocks}, "docx_write")
        text = _docx_text(path)
        check("docx truncation is flagged on the stamp", stamp.truncated)
        check("the docx states its bound INSIDE the document", "TRUNCATED" in text, text[:200])
        check("the bound value is in the document", "3" in text, text[:200])

        from pptx import Presentation

        slides = [{"layout": "bullets", "title": f"S{i}", "bullets": ["x"], "detail_ref": "d"}
                  for i in range(10)]
        p2, _c2, s2, _l2 = _build(pptx_writer.build, {"title": "T", "slides": slides}, "pptx_write")
        check("pptx truncation is flagged", s2.truncated)
        prs = Presentation(str(p2))
        deck_text = "\n".join(sh.text_frame.text for s in prs.slides
                              for sh in s.shapes if sh.has_text_frame)
        check("the deck states its bound on a slide", "TRUNCATED" in deck_text, deck_text[:200])

        import openpyxl

        rows = [[{"v": i, "src": "t"}] for i in range(10)]
        p3, _c3, s3, _l3 = _build(xlsx_writer.build,
                                  {"sheets": [{"name": "S", "columns": ["N"], "rows": rows}]},
                                  "xlsx_write")
        check("xlsx truncation is flagged", s3.truncated)
        ws = openpyxl.load_workbook(str(p3))["S"]
        sheet_text = "\n".join(str(ws.cell(row=r, column=1).value or "") for r in range(1, 6))
        check("the sheet states its bound", "TRUNCATED" in sheet_text, sheet_text[:200])
    finally:
        for k in ("DOCUMENT_MAX_BLOCKS", "DOCUMENT_MAX_SLIDES", "DOCUMENT_MAX_ROWS"):
            os.environ.pop(k, None)
        cleanup(d)


def test_disagreeing_sources_are_both_rendered():
    """FR-027b / SC-027 — no winner is picked and neither value is dropped."""
    d = sandbox()
    try:
        payload = {"title": "T", "blocks": [
            {"type": "figure", "label": "Hostname",
             "value": {"v": "fgt-01", "src": "fgt_system_status"}},
            {"type": "figure", "label": "Hostname",
             "value": {"v": "fgt-01-alt", "src": "netbox_get_device"}},
        ]}
        path, caveats, _s, _l = _build(docx_writer.build, payload, "docx_write")
        text = _docx_text(path)
        check("the first value is present", "fgt-01" in text)
        check("the second value is present", "fgt-01-alt" in text,
              "a value was dropped — the document asserts a certainty the data lacks")
        check("the disagreement is called out", "Sources disagree" in text, text[:200])
        check("both sources are named", "fgt_system_status" in text and "netbox_get_device" in text)
        check("a caveat reports it", any("disagree" in c for c in caveats), str(caveats))
        check("NetClaw says it did not reconcile", "not reconciled" in text.lower(), text[:300])
    finally:
        cleanup(d)


def test_prose_asserting_a_figure_is_flagged():
    d = sandbox()
    try:
        payload = {"title": "T", "blocks": [
            {"type": "paragraph", "text": "We observed 12 interfaces."},
            {"type": "paragraph", "text": "Captured 2026-08-03T14:02:00Z for CHG0012345 on 10.0.0.1 running v7.6.7, see Section 2."},
        ]}
        _p, caveats, _s, _l = _build(docx_writer.build, payload, "docx_write")
        flagged = [c for c in caveats if "prose asserts" in c]
        check("a bare figure in prose is flagged", len(flagged) == 1, str(caveats))
        if flagged:
            check("the caveat names the block index", "blocks[0]" in flagged[0], flagged[0])
        check(
            "dates, tickets, IPs, versions and ordinals are NOT flagged",
            not any("blocks[1]" in c for c in caveats),
            "the lint is noisy enough to be ignored, which makes it useless",
        )
    finally:
        cleanup(d)


TESTS = [
    test_gaps_render_explicitly_in_the_document,
    test_a_gapped_document_cannot_report_as_clean,
    test_failed_devices_are_rows_not_omissions,
    test_admin_and_oper_state_stay_separate,
    test_all_three_office_formats_refuse_a_template,
    test_bounds_are_stated_inside_the_document,
    test_disagreeing_sources_are_both_rendered,
    test_prose_asserting_a_figure_is_flagged,
]

if __name__ == "__main__":
    raise SystemExit(run(TESTS, "no-fabrication"))
