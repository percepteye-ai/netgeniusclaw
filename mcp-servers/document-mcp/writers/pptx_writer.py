"""Decks. Spec 082, US3.

Sources go ON the slide, in a shape, at the bottom.

python-pptx exposes `notes_slide`, and putting attribution there is the obvious move
and the wrong one: speaker notes are invisible in presentation mode and absent from a
default print or PDF export. That is the same class of hidden mechanism as a cell
comment, and FR-008a rules it out as the mechanism. Notes are written too, but only
additively.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from guards import apply_bound, max_slides, reject_template, resolve_embedded_image
from outcomes import (
    Outcome,
    TaggedValue,
    ValueRefused,
    find_disagreements,
    parse_tagged,
    render_source,
    render_tagged,
)
from provenance import DocumentStamp, SourceLedger
from sanitize import plain_text

MUTED = RGBColor(0x60, 0x60, 0x60)
GAP = RGBColor(0x99, 0x00, 0x00)

LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6


def build(payload: dict, stamp: DocumentStamp, ledger: SourceLedger) -> tuple[bytes, list[str]]:
    reject_template(payload)

    slides = payload.get("slides") or []
    if not isinstance(slides, list):
        raise ValueRefused("'slides' must be a list", Outcome.REFUSED_UNTYPED)

    slides, truncated = apply_bound(slides, max_slides())
    if truncated:
        stamp.truncated = True
        stamp.bound_applied = max_slides()
        stamp.bound_kind = "slides"

    caveats: list[str] = []
    labelled: list[tuple[str, TaggedValue]] = []

    prs = Presentation()
    _title_slide(prs, plain_text(payload.get("title") or "NetClaw findings"), stamp)

    for index, slide_spec in enumerate(slides):
        if not isinstance(slide_spec, dict):
            raise ValueRefused(f"slides[{index}] must be an object", Outcome.REFUSED_UNTYPED)
        path = f"slides[{index}]"
        layout = slide_spec.get("layout", "bullets")
        title = plain_text(slide_spec.get("title", ""))
        source_bits: list[str] = []

        if layout == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
            slide.shapes.title.text = title
            image = slide_spec.get("image") or {}
            resolved = resolve_embedded_image(
                image.get("path", ""), image.get("src", ""), f"{path}.image"
            )
            slide.shapes.add_picture(
                str(resolved), Inches(0.8), Inches(1.6), width=Inches(8.4)
            )
            source_bits.append(f"diagram produced by {plain_text(image['src'])}")

        elif layout == "figure":
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
            slide.shapes.title.text = title
            body = slide.placeholders[1].text_frame
            body.clear()
            figures = slide_spec.get("figures") or []
            for f_index, fig in enumerate(figures):
                label = plain_text(fig.get("label", ""))
                tv = ledger.record(parse_tagged(fig.get("value"), f"{path}.figures[{f_index}].value"))
                labelled.append((label, tv))
                para = body.paragraphs[0] if f_index == 0 else body.add_paragraph()
                para.text = f"{label}: {plain_text(render_tagged(tv))}"
                if tv.is_gap:
                    for run in para.runs:
                        run.font.color.rgb = GAP
                        run.font.bold = True
                bit = render_source(tv)
                if tv.kind == "value" and tv.as_of:
                    bit += f" (as of {tv.as_of})"
                if bit not in source_bits:
                    source_bits.append(bit)

        else:  # bullets / title
            slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
            slide.shapes.title.text = title
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = [plain_text(b) for b in (slide_spec.get("bullets") or [])]
            for b_index, bullet in enumerate(bullets):
                para = body.paragraphs[0] if b_index == 0 else body.add_paragraph()
                para.text = bullet
            detail_ref = slide_spec.get("detail_ref")
            if bullets and not detail_ref:
                caveats.append(
                    f"{path}: summary slide {title!r} has no 'detail_ref'. A summary "
                    f"claim should be traceable to a detail slide or the Sources slide "
                    f"rather than asserted bare (FR-005)"
                )
            if detail_ref:
                source_bits.append(f"detail: {plain_text(detail_ref)}")

        # The visible source line — a real shape, not speaker notes (FR-008a).
        _source_line(slide, prs, source_bits, stamp)
        # Additive only.
        slide.notes_slide.notes_text_frame.text = (
            (" · ".join(source_bits) if source_bits else "no figures on this slide")
            + f"\n{stamp.footer_text()}"
        )

    for dis in find_disagreements(labelled):
        caveats.append(dis.caveat())
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
        slide.shapes.title.text = f"Sources disagree — {dis.label}"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, v in enumerate(dis.values):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"{render_tagged(v)}  —  {render_source(v)}"
        tail = tf.add_paragraph()
        tail.text = "NetClaw has not reconciled these. Both are shown as reported."
        _source_line(slide, prs, ["values shown as reported by each source"], stamp)

    _sources_slide(prs, ledger, stamp)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), caveats


def _title_slide(prs: Presentation, title: str, stamp: DocumentStamp) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.shapes.title.text = title
    try:
        slide.placeholders[1].text = stamp.footer_text()
    except (KeyError, IndexError):
        _source_line(slide, prs, [], stamp)
    if stamp.truncated:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), prs.slide_width - Inches(1.0), Inches(0.6))
        para = box.text_frame.paragraphs[0]
        para.text = stamp.truncation_text()
        para.runs[0].font.size = Pt(11)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = GAP


def _source_line(slide, prs: Presentation, bits: list[str], stamp: DocumentStamp) -> None:
    """A visible textbox along the bottom of the slide. This is the mechanism —
    speaker notes are not."""
    text = " · ".join(b for b in bits if b)
    text = f"Source: {text}" if text else "No figures on this slide."
    height = Inches(0.4)
    top = prs.slide_height - height - Emu(45720)
    box = slide.shapes.add_textbox(Inches(0.4), top, prs.slide_width - Inches(0.8), height)
    para = box.text_frame.paragraphs[0]
    para.text = f"{text}   |   {stamp.generated_by} · {stamp.generated_at}"
    run = para.runs[0]
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED
    run.font.italic = True


def _sources_slide(prs: Presentation, ledger: SourceLedger, stamp: DocumentStamp) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    slide.shapes.title.text = "Sources"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    rows = ledger.as_rows()
    if not rows:
        tf.paragraphs[0].text = "No sources were consulted for this deck."
    for i, row in enumerate(rows):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        src, device, as_of, count, status = row
        where = f" · {device}" if device and device != "—" else ""
        para.text = f"{src}{where} — as of {as_of} — {count} element(s) — {status}"
        para.runs[0].font.size = Pt(14)
    _source_line(slide, prs, ["this slide lists every source consulted"], stamp)
